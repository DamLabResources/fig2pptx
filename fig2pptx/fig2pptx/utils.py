import pptx
import numpy as np


def bbox_to_pptx(left, bottom, width, height):
    """ Convert matplotlib bounding box format to pptx format
    Parameters
    ----------
    bottom : float
    left : float
    width : float
    height : float

    Returns
    -------
    left, top, width, height


    """

    return left, bottom-height, width, height


def pptx_to_bbox(left, top, width, height):
    """ Convert matplotlib bounding box format to pptx format
    Parameters
    ----------
    left : float
    top : float
    width : float
    height : float

    Returns
    -------
    bottom, left, width, height


    """

    return top-height, left, width, height


def corners_to_pptx(corners):

    top = corners[:,1].min()
    right = corners[:,0].max()
    left = corners[:,0].min()
    bottom = corners[:,1].max()

    height = bottom-top
    width = right-left

    return left, top, width, height


def transform_corners(mpl_corners, ref_corners = None, ref_shape = None):
    """
    Parameters
    ----------
    mpl_corners : np.array
    ref_corners : np.array
    ref_shape : np.array

    Returns
    -------
    np.array

    """

    if (ref_shape is None) and (ref_corners is None):
        raise ValueError('Must provide either a reference shape or reference corners')

    if ref_shape is not None:
        ref_corners = get_corners(ref_shape)

    new_corners = []

    for n in range(4):
        new_corners.append(transform_point(mpl_corners[n, :], ref_corners))

    return np.array(new_corners)



def transform_point(mpl_point, ref_corners, unit = pptx.util.Inches):
    """
    Parameters
    ----------
    mpl_point : np.array
    ref_shape : pptx.shapes.base.BaseShape

    Returns
    -------

    """

    left, bottom = ref_corners[:,0].min(), ref_corners[:,1].max()

    return np.array([left+unit(mpl_point[0]), bottom-unit(mpl_point[1])])


def get_corners(shape):
    """ Returns array of corners in pptx-space

    Parameters
    ----------
    shape : pptx.shapes.base.BaseShape

    Returns
    -------
    np.array

    """

    c1 = (shape.left, shape.top)
    c2 = (shape.left, shape.top + shape.height)
    c3 = (shape.left + shape.width, shape.top)
    c4 = (shape.left + shape.width, shape.top + shape.height)

    return np.array([c1, c2, c3, c4])


def find_textbox(slide, text = None, within = None):
    """ Finds textboxes on slides by parameters
    Parameters
    ----------
    slide : pptx.slide.Slide
    text : str
    within : np.array

    Returns
    -------
    pptx.shapes.autoshape.Shape

    """

    for shape in slide.shapes:
        if isinstance(shape, pptx.shapes.autoshape.Shape):
            if (text is not None) & (shape.text == text):
                return shape


def make_blank_slide():
    """ Makes new presentation with a single blank slide
    """

    prs = pptx.Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    return prs, slide

def determine_corners(fig, left_top):
    """

    Parameters
    ----------
    fig
    left_top

    Returns
    -------

    """

    wi, hi = fig.get_size_inches()
    hi = pptx.util.Inches(hi)
    wi = pptx.util.Inches(wi)

    left, top = left_top

    c1 = (left, top)
    c2 = (left, top + hi)
    c3 = (left + wi, top)
    c4 = (left + wi, top + hi)

    return np.array([c1, c2, c3, c4])
