import pptx
import matplotlib
from matplotlib.axes._subplots import _subplot_classes
from tempfile import NamedTemporaryFile
from . import utils


class PowerPointSlide(object):

    def __init__(self, pptx_slide):
        """
        Parameters
        ----------
        pptx_slide : pptx.slide.Slide
            The slide to draw onto.

        Returns
        -------

        """

        self.pptx_slide = pptx_slide
        self.figures = []

    def add_figure(self, power_point_box):
        """
        Parameters
        ----------
        power_point_box : PowerPointBox

        Returns
        -------
        None
        """
        self.figures.append(power_point_box)

    def render(self):

        for box in self.figures:
            box.render(self.pptx_slide)



class PowerPointBox(object):

    def __init__(self, mpl_object, parent=None, pptx_corners=None):
        """
        Parameters
        ----------
        mpl_object
            Any matplotlib object

        Returns
        -------

        """

        self.parent = parent
        self.mpl_object = mpl_object
        self.children = []
        self.position = None

        self.pptx_corners = pptx_corners

        for child in mpl_object.get_children():
            cls = PowerPointBox.detect_type(child)
            self.children.append(cls(child, parent=self))


    def __iter__(self):
        """  Depth-first traversal
        Returns
        -------

        """

        yield self
        for child in self.children:
            for node in child:
                yield node

    def get_top_fig(self):

        obj = self
        while obj.parent is not None:
            obj = obj.parent

        return obj

    @staticmethod
    def detect_type(mpl_object):

        for cls in PowerPointBox.__subclasses__():
            if cls.detect_type(mpl_object):
                return cls
        return PowerPointBox

    @property
    def figure_position(self):
        """In matplotlib corners format"""
        return None

    def slide_position(self):
        """In pptx (top, left, width, height) format using proper units!"""

        return utils.corners_to_pptx(self.pptx_corners)



    def _render(self, slide):

        if self.parent is None:
            # This is the head figure object
            with NamedTemporaryFile(suffix='.png') as handle:

                fig = self.mpl_object
                fig.savefig(handle.name, dpi=300, transparent=True)
                left, top, width, height = self.slide_position()
                pic = slide.shapes.add_picture(handle.name, left, top,
                                               height = height, width = width)

    def render(self, slide):

        self._render(slide)

        for child in self.children:
            child.render(slide)


class SubplotAxisBox(PowerPointBox):

    @staticmethod
    def detect_type(mpl_object):
        return isinstance(mpl_object, tuple(_subplot_classes))


    def _crop(self, img_box):
        """
        Parameters
        ----------
        img_box : pptx.shapes.picture.Picture

        Returns
        -------
        None

        """

        raise NotImplementedError

        # corners = self.mpl_object.get_position().corners()

        # img_box.crop_left(corners[:,0].min())
        # img_box.crop_right(1-corners[:,0].max())
        # img_box.crop_bottom(corners[:, 1].min())
        # img_box.crop_top(1-corners[:, 1].max())

    def _render(self, slide):

        pass

    def _future_render(self, slide):
        """ This should work when pptx implements cropping
        Parameters
        ----------
        slide : pptx.slide.Slide

        Returns
        -------

        """

        with NamedTemporaryFile(suffix='.png') as handle:
            head = self.get_top_fig()

            fig = head.mpl_object
            fig.savefig(handle.name, dpi=300, transparent=True)

            loc = head.slide_position()
            print(handle.name)
            pic = slide.shapes.add_picture(handle.name, *loc)
            self._crop(pic)


class TextBox(PowerPointBox):

    @staticmethod
    def detect_type(mpl_object):
        return isinstance(mpl_object, matplotlib.text.Text)


    @property
    def figure_position(self):

        head = self.get_top_fig()
        renderer = head.mpl_object.canvas.get_renderer()

        bbox = self.mpl_object.get_window_extent(renderer)
        bbox_trans = bbox.inverse_transformed(head.mpl_object.transFigure)
        return bbox_trans.corners()

    def slide_position(self):

        nudge = pptx.util.Inches(0.15) # determined by inspection

        head = self.get_top_fig()
        wi, hi = head.mpl_object.get_size_inches()

        box_corners = self.figure_position
        box_corners[:,0] *= wi
        box_corners[:,1] *= hi

        slide_corners = utils.transform_corners(box_corners, ref_corners = head.pptx_corners)
        slide_corners[:, 1] -= nudge

        return utils.corners_to_pptx(slide_corners)


    def _render(self, slide):
        """

        Parameters
        ----------
        slide : pptx.slide.Slide

        Returns
        -------

        """

        if (len(self.mpl_object.get_text())>0) & (not self.mpl_object.stale):

            position = self.slide_position()
            tbox = slide.shapes.add_textbox(*position)
            text_frame = tbox.text_frame
            text_frame.text = self.mpl_object.get_text()









