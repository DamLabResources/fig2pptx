import pytest
from fig2pptx import utils
from tempfile import NamedTemporaryFile
import pptx
import numpy as np

# class TestConvertBoxes(object):
#
#     def test_bbox_to_pptx(self):
#
#         # bottom, left, width, height
#         res = utils.bbox_to_pptx(5, 2, 4, 3)
#         assert res == (8, 2, 4, 3)
#
#
#     def test_pptx_to_bbox(self):
#
#         # top, left, width, height
#         res = utils.pptx_to_bbox(5, 2, 4, 3)
#         assert res == (2, 2, 4, 3)


def make_blank_slide():
    """ Makes new presentation with a single blank slide
    """

    prs = pptx.Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    return prs, slide


class TestFindObject(object):


    def test_find_textbox_by_text(self):

        with NamedTemporaryFile(suffix = '.pptx') as handle:

            prs, slide = make_blank_slide()
            position = (pptx.util.Inches(1),
                         pptx.util.Inches(2),
                         pptx.util.Inches(3),
                         pptx.util.Inches(4))

            tbox = slide.shapes.add_textbox(*position)
            tbox.text_frame.text = 'This is a textbox'

            prs.save(handle.name)

            nprs = pptx.Presentation(handle.name)
            slide = nprs.slides[0]

            new_tbox = utils.find_textbox(slide, text = 'This is a textbox')
            new_position = (new_tbox.left, new_tbox.top,
                            new_tbox.width, new_tbox.height)

            assert new_tbox.text_frame.text == 'This is a textbox'
            assert position == new_position




class TestConvertCorners(object):

    def test_get_corners(self):

        prs, slide = make_blank_slide()
        position = (pptx.util.Inches(1),
                     pptx.util.Inches(2),
                     pptx.util.Inches(3),
                     pptx.util.Inches(4))

        tbox = slide.shapes.add_textbox(*position)
        tbox.text_frame.text = 'This is a textbox'

        corners = utils.get_corners(tbox)
        tc1 = (pptx.util.Inches(1), pptx.util.Inches(2))
        tc2 = (pptx.util.Inches(1), pptx.util.Inches(2+4))
        tc3 = (pptx.util.Inches(1+3), pptx.util.Inches(2))
        tc4 = (pptx.util.Inches(1+3), pptx.util.Inches(2+4))

        true_corners = np.array([tc1, tc2, tc3, tc4])

        np.testing.assert_array_equal(corners, true_corners)


    def test_transform_point(self):

        prs, slide = make_blank_slide()
        position = (pptx.util.Inches(1),
                     pptx.util.Inches(2),
                     pptx.util.Inches(3),
                     pptx.util.Inches(4))

        tbox = slide.shapes.add_textbox(*position)
        tbox.text_frame.text = 'This is a textbox'

        mpl_point = np.array([2.1, 1.1])

        new_point = utils.transform_point(mpl_point, tbox)

        cor_point = np.array([pptx.util.Inches(1+2.1), pptx.util.Inches(6-1.1)])
        np.testing.assert_array_equal(new_point, cor_point)




